from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── palette ──────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0F, 0x17, 0x2A)   # deep navy
ACCENT   = RGBColor(0x38, 0xBD, 0xF8)   # sky-400
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
MUTED    = RGBColor(0x94, 0xA3, 0xB8)   # slate-400
CARD_BG  = RGBColor(0x1E, 0x29, 0x3B)   # slate-800

W = Inches(13.33)
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def bg(slide):
    shape = slide.shapes.add_shape(1, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG
    shape.line.fill.background()
    return shape


def rect(slide, x, y, w, h, color, radius=0):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def label(slide, text, x, y, w, h, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def accent_bar(slide, y, thick=Inches(0.04)):
    bar = slide.shapes.add_shape(1, Inches(0.6), y, Inches(2), thick)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


# ─────────────────────────────────────────────────────────────────────────────
# Slide 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
def slide_title(prs):
    s = blank_slide(prs)
    bg(s)

    # decorative glow circle
    circ = s.shapes.add_shape(9, Inches(9.5), Inches(-1), Inches(5), Inches(5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = RGBColor(0x0E, 0x4A, 0x6E)
    circ.line.fill.background()

    label(s, "Weather App",
          Inches(0.8), Inches(1.8), Inches(9), Inches(1.4),
          54, bold=True, color=WHITE)

    label(s, "Architecture  ·  Technologies  ·  Testing",
          Inches(0.8), Inches(3.1), Inches(9), Inches(0.6),
          22, color=ACCENT)

    label(s, "Next.js  ·  TanStack Query  ·  shadcn/ui  ·  Tailwind CSS  ·  TypeScript",
          Inches(0.8), Inches(3.9), Inches(10), Inches(0.5),
          14, color=MUTED)

    label(s, "2026", Inches(0.8), Inches(6.5), Inches(2), Inches(0.4),
          12, color=MUTED)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 2 — Tech Stack overview (cards)
# ─────────────────────────────────────────────────────────────────────────────
def slide_stack(prs):
    s = blank_slide(prs)
    bg(s)

    label(s, "Technology Stack", Inches(0.6), Inches(0.3), Inches(10), Inches(0.7),
          32, bold=True, color=WHITE)
    accent_bar(s, Inches(1.05))

    cards = [
        ("Next.js 16",         "App Router · React 19 · Server Components · Server Functions ('use server')"),
        ("TypeScript",          "Strict mode · end-to-end type safety from server actions to UI"),
        ("TanStack Query v5",  "Client-side server state · caching · background refetch · staleTime tuning"),
        ("shadcn/ui",           "Pre-built accessible component primitives · customised via className + cn()"),
        ("Tailwind CSS v4",    "Utility-first styling · PostCSS pipeline · no custom CSS primitives"),
        ("Vitest",              "Unit & integration tests · fetch mocking · runs in Node without a browser"),
    ]

    cols, rows = 3, 2
    cw, ch = Inches(3.9), Inches(2.1)
    gx, gy = Inches(0.4), Inches(1.35)
    pad = Inches(0.18)

    for i, (title, body) in enumerate(cards):
        col, row = i % cols, i // cols
        x = gx + col * (cw + Inches(0.17))
        y = gy + row * (ch + Inches(0.17))

        card = rect(s, x, y, cw, ch, CARD_BG)
        label(s, title, x + pad, y + pad, cw - 2*pad, Inches(0.4),
              14, bold=True, color=ACCENT)
        label(s, body,  x + pad, y + pad + Inches(0.42), cw - 2*pad, ch - pad - Inches(0.5),
              11, color=MUTED)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 3 — Architecture principles
# ─────────────────────────────────────────────────────────────────────────────
def slide_arch(prs):
    s = blank_slide(prs)
    bg(s)

    label(s, "Architecture Principles", Inches(0.6), Inches(0.3), Inches(10), Inches(0.7),
          32, bold=True, color=WHITE)
    accent_bar(s, Inches(1.05))

    principles = [
        ("Server-first rendering",
         "Pages are Server Components by default. 'use client' is added only for\n"
         "interactive islands — kept small and deep in the tree."),
        ("No API routes",
         "All backend logic lives in 'use server' functions. No route handlers,\n"
         "no fetch boilerplate — server functions are the only door to the backend."),
        ("No useEffect for data",
         "TanStack Query replaces every data-fetching useEffect.\n"
         "useEffect is reserved for non-React external systems only."),
        ("Single source of truth",
         "Server data lives in TanStack Query cache. URL/nav state in searchParams.\n"
         "Local UI state in useState. Nothing is duplicated."),
    ]

    for i, (title, body) in enumerate(principles):
        y = Inches(1.4) + i * Inches(1.4)
        num = rect(s, Inches(0.5), y + Inches(0.05), Inches(0.55), Inches(0.55),
                   ACCENT)
        label(s, str(i + 1), Inches(0.5), y + Inches(0.02), Inches(0.55), Inches(0.55),
              18, bold=True, color=BG, align=PP_ALIGN.CENTER)
        label(s, title, Inches(1.25), y, Inches(10), Inches(0.4),
              15, bold=True, color=WHITE)
        label(s, body,  Inches(1.25), y + Inches(0.42), Inches(11.2), Inches(0.8),
              12, color=MUTED)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 4 — Data flow diagram (text-based)
# ─────────────────────────────────────────────────────────────────────────────
def slide_flow(prs):
    s = blank_slide(prs)
    bg(s)

    label(s, "Data Flow", Inches(0.6), Inches(0.3), Inches(10), Inches(0.7),
          32, bold=True, color=WHITE)
    accent_bar(s, Inches(1.05))

    steps = [
        ("User Input",         "Search box (controlled useState) → debounced 300 ms"),
        ("TanStack Query",     "useQuery fires searchLocations() when debounced ≥ 2 chars\nqueryKey: ['geocode', debounced]  ·  staleTime: 5 min"),
        ("Server Function",    "'use server' geocode.ts  —  routes ZIP (3–5 digits) or city name\nCalls Census TIGERweb REST API · sanitises input · ranks results"),
        ("Response & Cache",   "Up to 8 GeocodeMatch objects returned and cached by TanStack Query\nDropdown rendered; user selects → weather query fires"),
        ("Weather Display",    "Coordinates passed to weather server action\nCurrent conditions · hourly forecast · 7-day daily forecast"),
    ]

    arrow_color = ACCENT
    bw, bh = Inches(2.0), Inches(0.72)
    bx = Inches(0.55)
    gap = Inches(0.25)

    for i, (step, detail) in enumerate(steps):
        y = Inches(1.3) + i * (bh + gap)

        box = rect(s, bx, y, bw, bh, CARD_BG)
        label(s, step, bx, y + Inches(0.14), bw, bh,
              11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

        label(s, detail, Inches(2.85), y + Inches(0.12), Inches(9.9), bh,
              11, color=MUTED)

        if i < len(steps) - 1:
            arr = s.shapes.add_shape(1, bx + bw/2 - Inches(0.03),
                                      y + bh, Inches(0.06), gap)
            arr.fill.solid()
            arr.fill.fore_color.rgb = ACCENT
            arr.line.fill.background()


# ─────────────────────────────────────────────────────────────────────────────
# Slide 5 — Testing
# ─────────────────────────────────────────────────────────────────────────────
def slide_testing(prs):
    s = blank_slide(prs)
    bg(s)

    label(s, "Testing Strategy", Inches(0.6), Inches(0.3), Inches(10), Inches(0.7),
          32, bold=True, color=WHITE)
    accent_bar(s, Inches(1.05))

    label(s, "Framework:  Vitest 4  (ESM-native, TypeScript, no browser required)",
          Inches(0.6), Inches(1.2), Inches(12), Inches(0.45), 14, color=WHITE)

    label(s, "Test file:  src/app/actions/__tests__/geocode.test.ts  —  10 tests, 3 suites",
          Inches(0.6), Inches(1.65), Inches(12), Inches(0.4), 13, color=MUTED)

    suites = [
        ("ZIP code search  (6 tests)",
         [
             "5-digit ZIP → ZCTA Census layer endpoint",
             "3-digit prefix → WHERE clause contains LIKE '900%'",
             "4-digit prefix routes correctly",
             "Enriched label: 'San Francisco, CA (94102)'",
             "Fallback label: 'ZIP 99999' when reverse lookup returns nothing",
             "API failure returns [] gracefully",
         ]),
        ("City routing  (2 tests)",
         [
             "City name does NOT hit the ZCTA layer",
             "1-char query returns [] without calling fetch",
         ]),
        ("Input sanitization  (2 tests)",
         [
             "Empty string returns [] without calling fetch",
             "SQL injection payload stripped from WHERE clause",
         ]),
    ]

    col_x = [Inches(0.55), Inches(4.7), Inches(9.0)]
    col_w = Inches(3.9)
    top_y = Inches(2.2)

    for ci, (suite_title, items) in enumerate(suites):
        x = col_x[ci]
        rect(s, x, top_y, col_w, Inches(4.6), CARD_BG)
        label(s, suite_title, x + Inches(0.15), top_y + Inches(0.12),
              col_w - Inches(0.3), Inches(0.5), 12, bold=True, color=ACCENT)
        for ri, item in enumerate(items):
            label(s, f"✓  {item}",
                  x + Inches(0.15),
                  top_y + Inches(0.65) + ri * Inches(0.58),
                  col_w - Inches(0.3), Inches(0.5),
                  11, color=MUTED)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 6 — Key decisions & trade-offs
# ─────────────────────────────────────────────────────────────────────────────
def slide_decisions(prs):
    s = blank_slide(prs)
    bg(s)

    label(s, "Key Decisions", Inches(0.6), Inches(0.3), Inches(10), Inches(0.7),
          32, bold=True, color=WHITE)
    accent_bar(s, Inches(1.05))

    decisions = [
        ("Server Functions over API routes",
         "Keeps the backend surface area tiny. No extra file, no HTTP verb to reason about.\n"
         "Direct TypeScript call — compiler catches mismatches end-to-end."),
        ("Census TIGERweb for geocoding",
         "Free, no API key, US-only scope fits the app. ZIP centroid lookup + reverse\n"
         "place-name enrichment avoids a paid geocoding service."),
        ("Debounced query (300 ms)",
         "Prevents a Census API call on every keystroke. staleTime 5 min means repeated\n"
         "searches for the same prefix are served from cache."),
        ("SQL-injection sanitization",
         "TIGERweb uses a WHERE clause in the URL. Input is stripped of non-alphanumeric\n"
         "chars and single-quotes are escaped before interpolation."),
        ("Vitest over Jest",
         "ESM-native — no babel transform needed for TypeScript or Next.js imports.\n"
         "'use server' directive is handled transparently; tests run in milliseconds."),
    ]

    for i, (title, body) in enumerate(decisions):
        y = Inches(1.3) + i * Inches(1.18)
        dot = rect(s, Inches(0.55), y + Inches(0.15), Inches(0.12), Inches(0.12), ACCENT)
        label(s, title, Inches(0.85), y, Inches(11.5), Inches(0.38),
              13, bold=True, color=WHITE)
        label(s, body, Inches(0.85), y + Inches(0.38), Inches(11.5), Inches(0.7),
              11, color=MUTED)


# ─────────────────────────────────────────────────────────────────────────────
# Build
# ─────────────────────────────────────────────────────────────────────────────
prs = new_prs()
slide_title(prs)
slide_stack(prs)
slide_arch(prs)
slide_flow(prs)
slide_testing(prs)
slide_decisions(prs)

out = "/Users/isaachennessey/weatherAPP/WeatherApp_Overview.pptx"
prs.save(out)
print(f"Saved → {out}")
